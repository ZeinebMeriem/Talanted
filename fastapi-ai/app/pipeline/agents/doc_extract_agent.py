from __future__ import annotations

import logging
import os
from io import BytesIO
from urllib.parse import urlparse

try:
    from minio import Minio  # type: ignore[import-not-found]
    from minio.error import S3Error  # type: ignore[import-not-found]
except Exception:  # noqa: BLE001
    Minio = None  # type: ignore[assignment]

    class S3Error(Exception):  # type: ignore[no-redef]
        pass

try:
    from pypdf import PdfReader  # type: ignore[import-not-found]
except Exception:  # noqa: BLE001
    PdfReader = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument  # type: ignore[import-not-found]
except Exception:  # noqa: BLE001
    DocxDocument = None  # type: ignore[assignment]

try:
    from pptx import Presentation as PptxPresentation  # type: ignore[import-not-found]
except Exception:  # noqa: BLE001
    PptxPresentation = None  # type: ignore[assignment]

from ..models import SourceItem, SourcePack

logger = logging.getLogger(__name__)

# Extensions treated as plain text (including code files)
_TEXT_EXTS: tuple[str, ...] = (
    ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
    # code / structured
    ".js", ".ts", ".tsx", ".jsx", ".html", ".htm", ".css",
    ".py", ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs",
    ".sql", ".sh", ".env.example",
    # diagrams (text-based)
    ".mmd", ".excalidraw",
)

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class DocExtractAgent:
    """Download uploaded files from MinIO and extract text.

    Supported formats:
      • Plain text / code : .txt .md .csv .json .xml .yaml .js .ts .html .css .py …
      • PDF (text-native)  : up to AI_EXTRACT_MAX_PDF_PAGES pages
      • Word (.docx)       : paragraphs + tables
      • PowerPoint (.pptx) : all slide text shapes
    """

    def run(self, pack: SourcePack) -> SourcePack:
        if Minio is None:
            return pack

        endpoint   = os.getenv("MINIO_ENDPOINT",   "").strip()
        access_key = os.getenv("MINIO_ACCESS_KEY",  "").strip()
        secret_key = os.getenv("MINIO_SECRET_KEY",  "").strip()
        bucket     = os.getenv("MINIO_BUCKET",      "").strip()

        if not (endpoint and access_key and secret_key and bucket):
            return pack

        parsed = urlparse(endpoint)
        secure = parsed.scheme == "https"
        host   = parsed.netloc or parsed.path

        client = Minio(host, access_key=access_key, secret_key=secret_key, secure=secure)

        max_bytes     = int(os.getenv("AI_EXTRACT_MAX_BYTES",    "5242880"))  # 5 MB
        max_chars     = int(os.getenv("AI_EXTRACT_MAX_CHARS",    "20000"))
        max_pdf_pages = int(os.getenv("AI_EXTRACT_MAX_PDF_PAGES", "15"))

        new_items: list[SourceItem] = []
        for it in pack.items:
            if it.kind != "file_ref":
                continue

            key  = (it.meta.get("minioPath") or it.content or "").strip()
            if not key:
                continue

            mime       = str(it.meta.get("mimeType") or "").lower()
            name       = str(it.meta.get("originalName") or key)
            name_lower = name.lower()

            # ── Classify ──────────────────────────────────────────────────
            is_textish = mime.startswith("text/") or any(name_lower.endswith(e) for e in _TEXT_EXTS)
            is_pdf  = (mime == "application/pdf")  or name_lower.endswith(".pdf")
            is_docx = (mime == _DOCX_MIME)         or name_lower.endswith(".docx")
            is_pptx = (mime == _PPTX_MIME)         or name_lower.endswith(".pptx")

            if not (is_textish or is_pdf or is_docx or is_pptx):
                continue  # images / other binary → handled by OcrAgent

            try:
                resp = client.get_object(bucket, key)
                try:
                    data = resp.read(max_bytes + 1)
                finally:
                    resp.close()
                    resp.release_conn()

                if len(data) > max_bytes:
                    logger.warning("DocExtractAgent: %s exceeds %d bytes, skipping", name, max_bytes)
                    continue

                text: str | None = None
                extracted_pages: int | None = None
                extract_method  = "text"

                # ── Plain text / code ──────────────────────────────────
                if is_textish:
                    text = data.decode("utf-8-sig", errors="replace")
                    extract_method = "text"

                # ── PDF (text-native) ──────────────────────────────────
                elif is_pdf:
                    extract_method = "pdf_text"
                    if PdfReader is None:
                        continue
                    reader = PdfReader(BytesIO(data))
                    parts: list[str] = []
                    for i, page in enumerate(reader.pages):
                        if i >= max_pdf_pages:
                            break
                        try:
                            parts.append(page.extract_text() or "")
                        except Exception:
                            parts.append("")
                    extracted_pages = min(len(reader.pages), max_pdf_pages)
                    text = "\n".join(p for p in parts if p)

                # ── Word (.docx) ───────────────────────────────────────
                elif is_docx:
                    extract_method = "docx"
                    if DocxDocument is None:
                        logger.warning("DocExtractAgent: python-docx not available, skipping %s", name)
                        continue
                    doc = DocxDocument(BytesIO(data))
                    lines: list[str] = []
                    # Paragraphs (headings + body)
                    for para in doc.paragraphs:
                        stripped = para.text.strip()
                        if stripped:
                            lines.append(stripped)
                    # Tables
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                lines.append(row_text)
                    text = "\n".join(lines)

                # ── PowerPoint (.pptx) ─────────────────────────────────
                elif is_pptx:
                    extract_method = "pptx"
                    if PptxPresentation is None:
                        logger.warning("DocExtractAgent: python-pptx not available, skipping %s", name)
                        continue
                    prs = PptxPresentation(BytesIO(data))
                    slide_texts: list[str] = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        shapes_text: list[str] = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                shapes_text.append(shape.text.strip())
                        if shapes_text:
                            slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(shapes_text))
                    text = "\n\n".join(slide_texts)

                if text is None:
                    continue

                cleaned = "\n".join(line.rstrip() for line in text.splitlines()).strip()
                extracted_chars = len(cleaned)

                if cleaned and len(cleaned) > max_chars:
                    cleaned = cleaned[:max_chars] + "\n..."

                new_items.append(SourceItem(
                    kind="file_text",
                    content=cleaned,
                    meta={
                        **it.meta,
                        "minioPath":      key,
                        "originalName":   name,
                        "mimeType":       mime,
                        "extract":        extract_method,
                        "extractedChars": extracted_chars,
                        "extractedPages": extracted_pages,
                        "extractEmpty":   (extracted_chars == 0),
                    },
                ))

            except S3Error as e:
                logger.warning("DocExtractAgent: S3 error for %s: %s", key, e)
                continue
            except Exception as e:
                logger.warning("DocExtractAgent: error extracting %s: %s", key, e)
                continue

        logger.info(
            "DocExtractAgent: processed %d file_ref items → %d file_text items",
            sum(1 for i in pack.items if i.kind == "file_ref"),
            len(new_items),
        )
        for ni in new_items:
            logger.info("  extracted: %s [%s] (%d chars)",
                        ni.meta.get("originalName", "?"),
                        ni.meta.get("extract", "?"),
                        len(ni.content))

        if not new_items:
            return pack

        return SourcePack(items=[*pack.items, *new_items])
